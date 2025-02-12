from typing import Dict, List, Optional, Union, BinaryIO
from abc import ABC, abstractmethod
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import pandas as pd
import openpyxl
from PIL import Image
import pytesseract
import cv2
import numpy as np
import io
from .image_processor import ImageProcessor
from .table_processor import TableProcessor

class FileProcessor(ABC):
    """文件处理基类"""
    
    def __init__(self):
        self.image_processor = ImageProcessor()
        self.table_processor = TableProcessor()
    
    @abstractmethod
    async def extract_tables(self, file_content: BinaryIO) -> List[Dict]:
        """提取表格内容"""
        pass
    
    @abstractmethod
    async def extract_charts(self, file_content: BinaryIO) -> List[Dict]:
        """提取图表内容"""
        pass
    
    @abstractmethod
    async def extract_images(self, file_content: BinaryIO) -> List[bytes]:
        """提取图片内容"""
        pass

class PDFProcessor(FileProcessor):
    """PDF文件处理器"""
    
    async def extract_tables(self, file_content: BinaryIO) -> List[Dict]:
        try:
            doc = fitz.open(stream=file_content.read(), filetype="pdf")
            tables = []
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                # 使用PyMuPDF的表格检测功能
                tab = page.find_tables()
                if tab.tables:
                    for idx, table in enumerate(tab.tables):
                        # 转换表格数据为DataFrame
                        df = pd.DataFrame(table.extract())
                        # 使用TableProcessor处理表格数据
                        table_info = self.table_processor.parse_table(df)
                        tables.append({
                            "page": page_num + 1,
                            "index": idx + 1,
                            "data": table_info
                        })
            
            return tables
            
        except Exception as e:
            raise Exception(f"PDF table extraction error: {str(e)}")
    
    async def extract_charts(self, file_content: BinaryIO) -> List[Dict]:
        try:
            doc = fitz.open(stream=file_content.read(), filetype="pdf")
            charts = []
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                # 提取页面中的图片
                image_list = page.get_images(full=True)
                
                for img_idx, img in enumerate(image_list):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    # 使用ImageProcessor分析图片
                    image_analysis = await self.image_processor.analyze_image(image_bytes)
                    
                    # 如果是图表类型，提取数据
                    if image_analysis["type"] == "chart":
                        chart_data = await self.image_processor.extract_chart_data(
                            Image.open(io.BytesIO(image_bytes))
                        )
                        charts.append({
                            "page": page_num + 1,
                            "index": img_idx + 1,
                            "analysis": image_analysis,
                            "data": chart_data
                        })
            
            return charts
            
        except Exception as e:
            raise Exception(f"PDF chart extraction error: {str(e)}")
    
    async def extract_images(self, file_content: BinaryIO) -> List[bytes]:
        try:
            doc = fitz.open(stream=file_content.read(), filetype="pdf")
            images = []
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                image_list = page.get_images(full=True)
                
                for img in image_list:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    images.append(base_image["image"])
            
            return images
            
        except Exception as e:
            raise Exception(f"PDF image extraction error: {str(e)}")

class WordProcessor(FileProcessor):
    """Word文件处理器"""
    
    async def extract_tables(self, file_content: BinaryIO) -> List[Dict]:
        try:
            doc = Document(file_content)
            tables = []
            
            for idx, table in enumerate(doc.tables):
                # 提取表格数据
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                
                # 转换为DataFrame并使用TableProcessor处理
                df = pd.DataFrame(table_data[1:], columns=table_data[0])
                table_info = self.table_processor.parse_table(df)
                
                tables.append({
                    "index": idx + 1,
                    "data": table_info
                })
            
            return tables
            
        except Exception as e:
            raise Exception(f"Word table extraction error: {str(e)}")
    
    async def extract_charts(self, file_content: BinaryIO) -> List[Dict]:
        try:
            doc = Document(file_content)
            charts = []
            
            for shape in doc.inline_shapes:
                if shape.type == 3:  # Chart type
                    # 提取图表图像
                    image_bytes = shape._inline.graphic.graphicData.pic.blipFill.blip.embed
                    
                    # 使用ImageProcessor分析图表
                    image_analysis = await self.image_processor.analyze_image(image_bytes)
                    
                    if image_analysis["type"] == "chart":
                        chart_data = await self.image_processor.extract_chart_data(
                            Image.open(io.BytesIO(image_bytes))
                        )
                        charts.append({
                            "index": len(charts) + 1,
                            "analysis": image_analysis,
                            "data": chart_data
                        })
            
            return charts
            
        except Exception as e:
            raise Exception(f"Word chart extraction error: {str(e)}")
    
    async def extract_images(self, file_content: BinaryIO) -> List[bytes]:
        try:
            doc = Document(file_content)
            images = []
            
            for shape in doc.inline_shapes:
                if shape.type == 3 or shape.type == 4:  # Picture or Chart
                    image_bytes = shape._inline.graphic.graphicData.pic.blipFill.blip.embed
                    images.append(image_bytes)
            
            return images
            
        except Exception as e:
            raise Exception(f"Word image extraction error: {str(e)}")

class PowerPointProcessor(FileProcessor):
    """PowerPoint文件处理器"""
    
    async def extract_tables(self, file_content: BinaryIO) -> List[Dict]:
        try:
            prs = Presentation(file_content)
            tables = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_table:
                        table_data = []
                        for row in shape.table.rows:
                            row_data = [cell.text.strip() for cell in row.cells]
                            table_data.append(row_data)
                        
                        # 转换为DataFrame并使用TableProcessor处理
                        df = pd.DataFrame(table_data[1:], columns=table_data[0])
                        table_info = self.table_processor.parse_table(df)
                        
                        tables.append({
                            "slide": slide.slide_id,
                            "data": table_info
                        })
            
            return tables
            
        except Exception as e:
            raise Exception(f"PowerPoint table extraction error: {str(e)}")
    
    async def extract_charts(self, file_content: BinaryIO) -> List[Dict]:
        try:
            prs = Presentation(file_content)
            charts = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_chart:
                        # 提取图表数据
                        chart = shape.chart
                        chart_data = {
                            "type": str(chart.chart_type),
                            "series": []
                        }
                        
                        # 提取系列数据
                        for series in chart.series:
                            series_data = {
                                "name": series.name,
                                "values": [v for v in series.values]
                            }
                            if hasattr(series, "categories"):
                                series_data["categories"] = [c for c in series.categories]
                            chart_data["series"].append(series_data)
                        
                        charts.append({
                            "slide": slide.slide_id,
                            "data": chart_data
                        })
            
            return charts
            
        except Exception as e:
            raise Exception(f"PowerPoint chart extraction error: {str(e)}")
    
    async def extract_images(self, file_content: BinaryIO) -> List[bytes]:
        try:
            prs = Presentation(file_content)
            images = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        # 提取图片数据
                        image_bytes = shape.image.blob
                        images.append(image_bytes)
            
            return images
            
        except Exception as e:
            raise Exception(f"PowerPoint image extraction error: {str(e)}")

class ExcelProcessor(FileProcessor):
    """Excel文件处理器"""
    
    async def extract_tables(self, file_content: BinaryIO) -> List[Dict]:
        try:
            wb = openpyxl.load_workbook(file_content)
            tables = []
            
            for sheet in wb.worksheets:
                # 获取工作表的数据范围
                data = []
                for row in sheet.iter_rows():
                    row_data = [cell.value for cell in row]
                    data.append(row_data)
                
                # 转换为DataFrame并使用TableProcessor处理
                df = pd.DataFrame(data[1:], columns=data[0])
                table_info = self.table_processor.parse_table(df)
                
                tables.append({
                    "sheet": sheet.title,
                    "data": table_info
                })
            
            return tables
            
        except Exception as e:
            raise Exception(f"Excel table extraction error: {str(e)}")
    
    async def extract_charts(self, file_content: BinaryIO) -> List[Dict]:
        try:
            wb = openpyxl.load_workbook(file_content)
            charts = []
            
            for sheet in wb.worksheets:
                for chart in sheet._charts:
                    chart_data = {
                        "type": chart.type,
                        "title": chart.title,
                        "series": []
                    }
                    
                    # 提取系列数据
                    for series in chart.series:
                        series_data = {
                            "title": series.title,
                            "values": [v for v in series.values]
                        }
                        if hasattr(series, "xvalues"):
                            series_data["categories"] = [x for x in series.xvalues]
                        chart_data["series"].append(series_data)
                    
                    charts.append({
                        "sheet": sheet.title,
                        "data": chart_data
                    })
            
            return charts
            
        except Exception as e:
            raise Exception(f"Excel chart extraction error: {str(e)}")
    
    async def extract_images(self, file_content: BinaryIO) -> List[bytes]:
        try:
            wb = openpyxl.load_workbook(file_content)
            images = []
            
            for sheet in wb.worksheets:
                for image in sheet._images:
                    images.append(image.ref)
            
            return images
            
        except Exception as e:
            raise Exception(f"Excel image extraction error: {str(e)}")

def get_file_processor(file_type: str) -> FileProcessor:
    """根据文件类型获取对应的处理器"""
    processors = {
        "pdf": PDFProcessor,
        "docx": WordProcessor,
        "pptx": PowerPointProcessor,
        "xlsx": ExcelProcessor
    }
    
    processor_class = processors.get(file_type.lower())
    if not processor_class:
        raise ValueError(f"Unsupported file type: {file_type}")
    
    return processor_class() 

# # 批量处理文件
# files = ["doc1.pdf", "doc2.docx", "presentation.pptx"]
# results = await extractor.batch_process_files(files)

# # 提取特定类型的内容
# charts = await extractor.extract_content_by_type("report.pdf", "charts")