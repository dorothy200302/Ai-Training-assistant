from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import List, Dict, Optional
import os
import re

class SlideTemplate:
    def __init__(
        self,
        layout_name: str,
        title_font_size: int = 32,
        body_font_size: int = 18,
        title_font_name: str = "Arial",
        body_font_name: str = "Arial"
    ):
        self.layout_name = layout_name
        self.title_font_size = title_font_size
        self.body_font_size = body_font_size
        self.title_font_name = title_font_name
        self.body_font_name = body_font_name

class PPTGenerator:
    """处理PPT生成的类"""
    
    def __init__(self, template_path: Optional[str] = None):
        self.template_path = template_path
        self.presentation = None
        self.slide_templates = {
            'title': SlideTemplate('Title Slide', 44, 32),
            'content': SlideTemplate('Content with Caption', 32, 18),
            'section': SlideTemplate('Section Header', 40, 24),
            'two_content': SlideTemplate('Two Content', 32, 18),
            'comparison': SlideTemplate('Comparison', 32, 18)
        }
        
        # 初始化演示文稿
        self._initialize_presentation()
    
    def _initialize_presentation(self):
        """初始化演示文稿"""
        if self.template_path and os.path.exists(self.template_path):
            self.presentation = Presentation(self.template_path)
        else:
            self.presentation = Presentation()
    
    def create_presentation(self, content: Dict[str, str]) -> None:
        """创建PPT演示文稿"""
        try:
            # 添加封面
            if 'title' in content:
                self.add_title_slide(
                    title=content['title'],
                    subtitle=content.get('subtitle', '')
                )
            
            # 添加目录
            if 'sections' in content:
                self.add_toc_slide(content['sections'])
            
            # 添加内容页
            for section in content.get('sections', []):
                self.add_section_slide(section['title'])
                
                # 处理章节内容
                if 'content' in section:
                    content_slides = self._split_content(section['content'])
                    for slide_content in content_slides:
                        self.add_content_slide(
                            title=section['title'],
                            content=slide_content
                        )
                
                # 处理子章节
                for subsection in section.get('subsections', []):
                    self.add_content_slide(
                        title=subsection['title'],
                        content=subsection.get('content', '')
                    )
            
        except Exception as e:
            raise Exception(f"Error creating presentation: {str(e)}")
    
    def add_title_slide(self, title: str, subtitle: str = '') -> None:
        """添加标题页"""
        slide_layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = title
        self._apply_text_style(
            title_shape.text_frame,
            self.slide_templates['title'].title_font_size,
            self.slide_templates['title'].title_font_name
        )
        
        # 设置副标题
        if subtitle and slide.placeholders[1]:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            self._apply_text_style(
                subtitle_shape.text_frame,
                self.slide_templates['title'].body_font_size,
                self.slide_templates['title'].body_font_name
            )
    
    def add_toc_slide(self, sections: List[Dict]) -> None:
        """添加目录页"""
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = "目录"
        self._apply_text_style(
            title_shape.text_frame,
            self.slide_templates['section'].title_font_size,
            self.slide_templates['section'].title_font_name
        )
        
        # 添加目录内容
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        for idx, section in enumerate(sections, 1):
            p = tf.add_paragraph()
            p.text = f"{idx}. {section['title']}"
            p.level = 0
            self._apply_text_style(
                tf,
                self.slide_templates['content'].body_font_size,
                self.slide_templates['content'].body_font_name
            )
            
            # 添加子章节
            for sub_idx, subsection in enumerate(section.get('subsections', []), 1):
                p = tf.add_paragraph()
                p.text = f"   {idx}.{sub_idx} {subsection['title']}"
                p.level = 1
    
    def add_section_slide(self, title: str) -> None:
        """添加章节页"""
        slide_layout = self.presentation.slide_layouts[2]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = title
        self._apply_text_style(
            title_shape.text_frame,
            self.slide_templates['section'].title_font_size,
            self.slide_templates['section'].title_font_name
        )
    
    def add_content_slide(self, title: str, content: str) -> None:
        """添加内容页"""
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = title
        self._apply_text_style(
            title_shape.text_frame,
            self.slide_templates['content'].title_font_size,
            self.slide_templates['content'].title_font_name
        )
        
        # 设置内容
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        # 处理内容中的项目符号
        for line in content.split('\n'):
            p = tf.add_paragraph()
            stripped_line = line.lstrip('- *')
            p.text = stripped_line
            p.level = 0 if line == stripped_line else 1
        
        self._apply_text_style(
            tf,
            self.slide_templates['content'].body_font_size,
            self.slide_templates['content'].body_font_name
        )
    
    def _split_content(self, content: str, max_chars_per_slide: int = 800) -> List[str]:
        """将内容分割成适合每页的大小"""
        if len(content) <= max_chars_per_slide:
            return [content]
        
        slides_content = []
        paragraphs = content.split('\n\n')
        current_slide = []
        current_length = 0
        
        for para in paragraphs:
            if current_length + len(para) > max_chars_per_slide:
                if current_slide:
                    slides_content.append('\n\n'.join(current_slide))
                    current_slide = []
                    current_length = 0
            
            current_slide.append(para)
            current_length += len(para)
        
        if current_slide:
            slides_content.append('\n\n'.join(current_slide))
        
        return slides_content
    
    def _apply_text_style(
        self,
        text_frame,
        font_size: int,
        font_name: str,
        color: tuple = (0, 0, 0)
    ) -> None:
        """应用文本样式"""
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.name = font_name
            paragraph.font.color.rgb = RGBColor(*color)
    
    def save(self, filename: str) -> None:
        """保存演示文稿"""
        try:
            self.presentation.save(filename)
        except Exception as e:
            raise Exception(f"Error saving presentation: {str(e)}") 